"""File loaders for Knowledge Base ingestion.

Each loader returns a list of `(text, piece_metadata)` tuples. `piece_metadata`
carries a location hint (page / slide / line_range) that the chunker attaches
to every derived chunk so citations can point back to the source region.
"""

from __future__ import annotations

from pathlib import Path
from typing import Callable, Dict, List, Tuple

TEXT_EXTS = {".txt", ".md", ".markdown"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}


MIME_BY_EXT = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppt": "application/vnd.ms-powerpoint",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".bmp": "image/bmp",
    ".webp": "image/webp",
}


def detect_mime(path: str) -> str:
    ext = Path(path).suffix.lower()
    return MIME_BY_EXT.get(ext, "application/octet-stream")


def is_image_mime(mime: str) -> bool:
    return mime.startswith("image/")


def load_pdf(path: str) -> List[Tuple[str, Dict]]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    out: List[Tuple[str, Dict]] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            out.append((text, {"page": i}))
    return out


def load_docx(path: str) -> List[Tuple[str, Dict]]:
    import docx2txt

    text = docx2txt.process(path) or ""
    if not text.strip():
        return []
    return [(text, {})]


def load_pptx(path: str) -> List[Tuple[str, Dict]]:
    from pptx import Presentation

    prs = Presentation(path)
    out: List[Tuple[str, Dict]] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text:
                            parts.append(run.text)
            elif hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        slide_text = "\n".join(parts).strip()
        if slide_text:
            out.append((slide_text, {"slide": i}))
    return out


def load_text(path: str) -> List[Tuple[str, Dict]]:
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    if not text.strip():
        return []
    return [(text, {})]


LOADERS: Dict[str, Callable[[str], List[Tuple[str, Dict]]]] = {
    "application/pdf": load_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": load_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": load_pptx,
    "text/plain": load_text,
    "text/markdown": load_text,
}


def load(path: str, mime: str | None = None) -> Tuple[str, List[Tuple[str, Dict]]]:
    """Dispatch to the right loader. Returns (mime, pieces)."""
    mime = mime or detect_mime(path)
    loader = LOADERS.get(mime)
    if loader is None:
        raise ValueError(f"Unsupported mime type for KB ingestion: {mime}")
    return mime, loader(path)
